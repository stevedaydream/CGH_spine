import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 1. Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Color Palette
    DARK_BLUE = RGBColor(13, 71, 161)   # #0D47A1
    LIGHT_BLUE = RGBColor(225, 245, 254) # #E1F5FE
    TEAL = RGBColor(0, 150, 136)        # #009688
    GREEN_BOX = RGBColor(46, 125, 50)   # #2E7D32
    RED_ACCENT = RGBColor(198, 40, 40)   # #C62828
    GRAY_BG = RGBColor(245, 245, 245)    # #F5F5F5
    GRAY_BORDER = RGBColor(200, 200, 200) # #C8C8C8
    TEXT_DARK = RGBColor(33, 33, 33)     # #212121
    TEXT_MUTED = RGBColor(117, 117, 117) # #757575
    WHITE = RGBColor(255, 255, 255)

    FONT_NAME = "Microsoft JhengHei"

    # Helper function to style shapes
    def style_shape(shape, fill_color, border_color=None):
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.color.rgb = fill_color

    # Helper function to add paragraph
    def add_para(tf, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, space_after=6):
        p = tf.add_paragraph() if tf.paragraphs and tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(space_after)
        return p

    # Helper function to add slide title
    def add_slide_title(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_NAME
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        return title_box

    # Helper function to add slide footer/page number
    def add_slide_footer(slide, page_num, total_pages=6):
        footer_box = slide.shapes.add_textbox(Inches(10.8), Inches(7.0), Inches(1.7), Inches(0.4))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = f"{page_num} / {total_pages}"
        p.font.name = FONT_NAME
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.RIGHT

    # Helper function to ensure white background
    def set_white_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # Slide 1: Cover
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide1)
    add_slide_footer(slide1, 1)

    # Add decorative background elements
    bottom_band = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
    style_shape(bottom_band, DARK_BLUE)

    # Add abstract spine icon
    cord = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.616), Inches(1.0), Inches(0.1), Inches(1.3))
    style_shape(cord, LIGHT_BLUE)

    for i in range(4):
        vert = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.266), Inches(1.05 + i*0.3), Inches(0.8), Inches(0.18))
        style_shape(vert, DARK_BLUE if i % 2 == 0 else TEAL)

    # Title & Subtitle text box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.333), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    add_para(tf1, "脊椎術後追蹤系統", size=44, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, space_after=12)
    add_para(tf1, "以數位工具實現主動式術後照護管理", size=20, bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER, space_after=10)

    # Department info text box
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    add_para(tf_info, "報告人：劉主任  ｜  指導單位：骨科部", size=15, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER, space_after=4)
    add_para(tf_info, "2026年6月", size=13, bold=False, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 2: Clinical Pain Points (Why Now)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide2)
    add_slide_title(slide2, "臨床痛點（Why Now）")
    add_slide_footer(slide2, 2)

    # Left column Header
    left_header = slide2.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5))
    tf_lh = left_header.text_frame
    add_para(tf_lh, "現況困境與照護瓶頸", size=18, bold=True, color=DARK_BLUE)

    pain_points = [
        "術後 12 週、17 個追蹤時間點，全靠人工無法落實",
        "電話追蹤：每位病患平均耗費護理師 [X] 分鐘/次",
        "資料無法跨病患比較，難以支撐臨床改善與研究發表"
    ]

    y_start = Inches(2.0)
    card_h = Inches(1.4)
    spacing = Inches(0.2)

    for i, pt in enumerate(pain_points):
        # Draw a card container
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_start + i * (card_h + spacing), Inches(5.5), card_h)
        style_shape(card, GRAY_BG, GRAY_BORDER)
        
        # Add a small red X icon
        x_icon = slide2.shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY, Inches(1.1), y_start + i * (card_h + spacing) + Inches(0.45), Inches(0.5), Inches(0.5))
        style_shape(x_icon, RED_ACCENT)
        
        # Add text
        tb = slide2.shapes.add_textbox(Inches(1.8), y_start + i * (card_h + spacing) + Inches(0.15), Inches(4.3), card_h - Inches(0.3))
        tf_tb = tb.text_frame
        tf_tb.word_wrap = True
        tf_tb.margin_left = tf_tb.margin_right = tf_tb.margin_top = tf_tb.margin_bottom = 0
        p = tf_tb.paragraphs[0]
        p.text = pt
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # Right column: Quote Box
    quote_box_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.6))
    style_shape(quote_box_bg, LIGHT_BLUE)

    # Large quotation marks
    quote_mark_start = slide2.shapes.add_textbox(Inches(7.3), Inches(1.8), Inches(1.0), Inches(0.8))
    tf_qms = quote_mark_start.text_frame
    add_para(tf_qms, "「", size=48, bold=True, color=DARK_BLUE)

    quote_text_box = slide2.shapes.add_textbox(Inches(7.6), Inches(2.6), Inches(4.3), Inches(2.4))
    tf_qt = quote_text_box.text_frame
    tf_qt.word_wrap = True
    add_para(tf_qt, "我們需要一個讓病患主動回報、\n讓醫護即時掌握、\n讓數據自動累積的系統。", size=22, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, space_after=12)

    quote_mark_end = slide2.shapes.add_textbox(Inches(11.0), Inches(4.8), Inches(1.0), Inches(0.8))
    tf_qme = quote_mark_end.text_frame
    add_para(tf_qme, "」", size=48, bold=True, color=DARK_BLUE, align=PP_ALIGN.RIGHT)


    # ==========================================
    # Slide 3: System Architecture (What)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide3)
    add_slide_title(slide3, "系統架構（What）")
    add_slide_footer(slide3, 3)

    # Layer 1: 病患端
    layer1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.4), Inches(10.33), Inches(1.2))
    style_shape(layer1, DARK_BLUE)
    tf_l1 = layer1.text_frame
    tf_l1.word_wrap = True
    add_para(tf_l1, "第一層：病患端（LINE Bot）", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=4)
    add_para(tf_l1, "零 App 安裝、零病患學習成本！Quick Reply 快速填答，保障極佳的依從性", size=13, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Arrow 1
    arrow1 = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.16), Inches(2.7), Inches(1.0), Inches(0.4))
    style_shape(arrow1, LIGHT_BLUE, DARK_BLUE)
    
    # Arrow 1 Label
    arrow1_lbl = slide3.shapes.add_textbox(Inches(7.3), Inches(2.7), Inches(4.0), Inches(0.4))
    tf_a1 = arrow1_lbl.text_frame
    add_para(tf_a1, "自動推播問卷 / Quick Reply 填答", size=12, bold=True, color=DARK_BLUE)

    # Layer 2: 後端
    layer2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.2), Inches(10.33), Inches(1.2))
    style_shape(layer2, GREEN_BOX)
    tf_l2 = layer2.text_frame
    tf_l2.word_wrap = True
    add_para(tf_l2, "第二層：後端 API（Google Apps Script）", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=4)
    add_para(tf_l2, "自動整合病患回報數據，進行 AI 輔助解析與疼痛惡化異常指標篩選警示", size=13, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow 2
    arrow2 = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.16), Inches(4.5), Inches(1.0), Inches(0.4))
    style_shape(arrow2, LIGHT_BLUE, DARK_BLUE)
    
    # Arrow 2 Label
    arrow2_lbl = slide3.shapes.add_textbox(Inches(7.3), Inches(4.5), Inches(4.0), Inches(0.4))
    tf_a2 = arrow2_lbl.text_frame
    add_para(tf_a2, "AI 輔助解析 / 異常主動警示", size=12, bold=True, color=DARK_BLUE)

    # Layer 3: 醫護端
    layer3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(10.33), Inches(1.2))
    style_shape(layer3, TEAL)
    tf_l3 = layer3.text_frame
    tf_l3.word_wrap = True
    add_para(tf_l3, "第三層：醫護端（網頁儀表板）", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=4)
    add_para(tf_l3, "個案管理儀表板，提供即時查閱、異常紅字警示確認、與一鍵匯出結構化臨床報告", size=13, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Highlight note at the bottom
    highlight_note = slide3.shapes.add_textbox(Inches(1.5), Inches(6.3), Inches(10.33), Inches(0.5))
    tf_hn = highlight_note.text_frame
    add_para(tf_hn, "★ 強調：零 App 安裝成本，病患零學習負擔，以現有 LINE 生態系實現無縫對接", size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)


    # ==========================================
    # Slide 4: System Screenshot / Live Demo (Demo)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide4)
    add_slide_title(slide4, "介面展示（Demo）")
    add_slide_footer(slide4, 4)

    # Left: Mock Browser with Web Screenshot
    # Draw Mock Browser Header
    browser_header = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(7.2), Inches(0.4))
    style_shape(browser_header, RGBColor(230, 230, 230))
    
    # Draw 3 window controls (circles)
    circle_colors = [RGBColor(255, 95, 86), RGBColor(255, 189, 46), RGBColor(39, 201, 63)]
    for idx, color in enumerate(circle_colors):
        circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95 + idx * 0.22), Inches(1.52), Inches(0.12), Inches(0.12))
        style_shape(circle, color)

    # Draw Address Bar
    address_bar = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(1.45), Inches(5.8), Inches(0.3))
    style_shape(address_bar, WHITE, GRAY_BORDER)
    tf_addr = address_bar.text_frame
    tf_addr.word_wrap = True
    tf_addr.margin_left = tf_addr.margin_right = Inches(0.1)
    tf_addr.margin_top = tf_addr.margin_bottom = Inches(0.02)
    p_addr = tf_addr.paragraphs[0]
    p_addr.text = "https://spine-cgh.pages.dev/demo"
    p_addr.font.name = FONT_NAME
    p_addr.font.size = Pt(9)
    p_addr.font.color.rgb = TEXT_MUTED
    p_addr.alignment = PP_ALIGN.LEFT

    # Insert Screenshot Image
    image_path = "spine_demo.png"
    if os.path.exists(image_path):
        pic = slide4.shapes.add_picture(image_path, Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.1))
        pic.click_action.hyperlink.address = "https://spine-cgh.pages.dev/demo"
    else:
        # Fallback card
        fallback = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.1))
        style_shape(fallback, GRAY_BG, GRAY_BORDER)
        tf_fb = fallback.text_frame
        tf_fb.word_wrap = True
        add_para(tf_fb, "\n未能載入 Demo 網頁截圖，請手動置入", size=18, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Under-screenshot instruction
    instruction_box = slide4.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(7.2), Inches(0.5))
    tf_inst = instruction_box.text_frame
    tf_inst.word_wrap = True
    add_para(tf_inst, "💡 提示：點擊上方畫面即可開啟線上 Demo 網頁進行互動體驗", size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Right: Description Card for the 3 Views
    desc_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.4), Inches(4.2), Inches(4.5))
    style_shape(desc_card, WHITE, DARK_BLUE)
    
    tf_desc = desc_card.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_right = Inches(0.2)
    tf_desc.margin_top = tf_desc.margin_bottom = Inches(0.15)
    
    add_para(tf_desc, "三大整合介面說明", size=18, bold=True, color=DARK_BLUE, space_after=12)
    
    add_para(tf_desc, "1. 護理師端（登記中心）", size=14, bold=True, color=TEXT_DARK, space_after=2)
    add_para(tf_desc, "搜尋病患、即時掌握術後天數、最近一次 VAS 疼痛指數與 LINE 綁定狀態。", size=12, bold=False, color=TEXT_MUTED, space_after=8)
    
    add_para(tf_desc, "2. 醫師端（管理後台）", size=14, bold=True, color=TEXT_DARK, space_after=2)
    add_para(tf_desc, "檢視追蹤完整度進度條、疼痛異常顏色警示（綠→紅）、AI 暫存待確認專區。", size=12, bold=False, color=TEXT_MUTED, space_after=8)
    
    add_para(tf_desc, "3. 病患端（LINE 問卷）", size=14, bold=True, color=TEXT_DARK, space_after=2)
    add_para(tf_desc, "直覺問卷與 Quick Reply 按鈕，降低長輩使用門檻，實現高回報率。", size=12, bold=False, color=TEXT_MUTED, space_after=12)
    
    # Highlights with red borders or bold red text
    add_para(tf_desc, "📌 重點標記：「異常 VAS 自動標紅」", size=13, bold=True, color=RED_ACCENT, space_after=2)
    add_para(tf_desc, "📌 重點標記：「追蹤完整度一目了然」", size=13, bold=True, color=TEAL, space_after=0)


    # ==========================================
    # Slide 5: Preliminary Results (So What)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide5)
    add_slide_title(slide5, "初步成果與預期效益（So What）")
    add_slide_footer(slide5, 5)

    cards_data = [
        {
            "title": "追蹤完整度",
            "highlight": "目標回報率由不規律提升至 ≥ 80%",
            "desc": "透過 LINE 自動化主動推播與零阻力的 Quick Reply 填答介面，大幅提升病患術後 12 週的填答依從性與資料完整度。",
            "bar_color": DARK_BLUE,
            "x": Inches(1.0), "y": Inches(1.6)
        },
        {
            "title": "護理追蹤人力",
            "highlight": "電話催填工作大幅減少",
            "desc": "系統自動推播與收集問卷，護理師僅需針對未填答或警示個案進行介入，預估可節省大量電話催填人力成本。",
            "bar_color": TEAL,
            "x": Inches(6.833), "y": Inches(1.6)
        },
        {
            "title": "研究資料庫",
            "highlight": "累積 12 週結構化資料",
            "desc": "自動累積標準化臨床結果指標，建立高價值研究資料庫，直接支援 VAS / ODI / MCID / PASS 等研究級統計分析。",
            "bar_color": GREEN_BOX,
            "x": Inches(1.0), "y": Inches(4.2)
        },
        {
            "title": "早期警示機制",
            "highlight": "疼痛惡化個案主動介入",
            "desc": "當病患回報評分異常（如疼痛急遽上升）時，系統發出即時警示，個管師可於下次門診回診前主動電話介入與關懷。",
            "bar_color": RED_ACCENT,
            "x": Inches(6.833), "y": Inches(4.2)
        }
    ]

    for card in cards_data:
        # Outer container
        container = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card["x"], card["y"], Inches(5.5), Inches(2.3))
        style_shape(container, WHITE, GRAY_BORDER)
        
        # Left colored accent bar
        accent_bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, card["x"], card["y"], Inches(0.15), Inches(2.3))
        style_shape(accent_bar, card["bar_color"])
        
        # Content text box
        content_box = slide5.shapes.add_textbox(card["x"] + Inches(0.3), card["y"] + Inches(0.15), Inches(5.0), Inches(2.0))
        tf_c = content_box.text_frame
        tf_c.word_wrap = True
        add_para(tf_c, card["title"], size=18, bold=True, color=DARK_BLUE, space_after=4)
        add_para(tf_c, card["highlight"], size=14, bold=True, color=card["bar_color"], space_after=8)
        add_para(tf_c, card["desc"], size=12, bold=False, color=TEXT_DARK, space_after=0)


    # ==========================================
    # Slide 6: Future Outlook
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide6)
    add_slide_title(slide6, "未來展望（針對部長）")
    add_slide_footer(slide6, 6)

    # Top Quote Box
    quote_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.4), Inches(10.33), Inches(1.1))
    style_shape(quote_box, LIGHT_BLUE)
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    add_para(tf_q, "「模模組化設計，一套邏輯，跨科複製」", size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Center Pilot Box
    center_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.416), Inches(2.9), Inches(2.5), Inches(1.0))
    style_shape(center_box, DARK_BLUE)
    tf_cb = center_box.text_frame
    tf_cb.word_wrap = True
    add_para(tf_cb, "脊椎骨科 Pilot", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=2)
    add_para(tf_cb, "（示範與標準流程建立）", size=11, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # 3 Expansion Target Cards
    targets = [
        {
            "title": "心臟外科",
            "subtitle": "CABG / 瓣膜置換術後\nHRQoL 追蹤 (SF-36、KCCQ)",
            "x": Inches(1.2), "y": Inches(4.5), "width": Inches(3.2), "height": Inches(1.2)
        },
        {
            "title": "一般外科",
            "subtitle": "癌症術後追蹤管理\n(自訂結構化收集機制)",
            "x": Inches(5.066), "y": Inches(4.5), "width": Inches(3.2), "height": Inches(1.2)
        },
        {
            "title": "骨科",
            "subtitle": "關節置換術後\n(KOOS、OHS 量表)",
            "x": Inches(8.933), "y": Inches(4.5), "width": Inches(3.2), "height": Inches(1.2)
        }
    ]

    for t in targets:
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t["x"], t["y"], t["width"], t["height"])
        style_shape(card, GRAY_BG, DARK_BLUE)
        tf_t = card.text_frame
        tf_t.word_wrap = True
        add_para(tf_t, t["title"], size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, space_after=4)
        add_para(tf_t, t["subtitle"], size=11, bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Draw arrows pointing down/outward
    # Arrow 1: to Left Card
    arrow_left = slide6.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.8), Inches(4.0), Inches(0.4), Inches(0.4))
    style_shape(arrow_left, TEAL)
    arrow_left.rotation = 45  # Points down-left (clockwise in PPTX)

    # Arrow 2: to Middle Card
    arrow_middle = slide6.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.466), Inches(4.0), Inches(0.4), Inches(0.4))
    style_shape(arrow_middle, TEAL)

    # Arrow 3: to Right Card
    arrow_right = slide6.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.133), Inches(4.0), Inches(0.4), Inches(0.4))
    style_shape(arrow_right, TEAL)
    arrow_right.rotation = 315 # Points down-right

    # Bottom Concluding Text Box
    conclusion_box = slide6.shapes.add_textbox(Inches(1.5), Inches(6.1), Inches(10.33), Inches(0.6))
    tf_c = conclusion_box.text_frame
    tf_c.word_wrap = True
    add_para(tf_c, "建議下一步：以骨科脊椎為 Pilot，建立院內數位術後追蹤標準流程", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Save presentation
    filename = "脊椎術後追蹤系統_劉主任簡報.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == "__main__":
    create_presentation()
